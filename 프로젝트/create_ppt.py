"""부동산 웹 플랫폼 제작 결과 발표 PPT 생성 스크립트"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUTPUT = "부동산_웹플랫폼_제작결과_발표.pptx"

# 색상 테마 (부동산/신뢰감 블루-그린)
PRIMARY = RGBColor(0x1A, 0x3A, 0x5C)
ACCENT = RGBColor(0x00, 0x96, 0x88)
LIGHT = RGBColor(0xF0, 0xF4, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x2C, 0x3E, 0x50)
SUBTEXT = RGBColor(0x5D, 0x6D, 0x7E)


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_header_bar(slide, title_text, prs):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.1)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.22), Inches(9), Inches(0.7))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "맑은 고딕"


def add_bullets(slide, items, left=0.6, top=1.5, width=8.8, height=5.5, font_size=18):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK_TEXT
        p.font.name = "맑은 고딕"
        p.space_after = Pt(10)
        p.level = 0


def add_two_column(slide, left_items, right_items, left_title="", right_title=""):
    # Left column
    if left_title:
        lt = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.2), Inches(0.5))
        lp = lt.text_frame.paragraphs[0]
        lp.text = left_title
        lp.font.size = Pt(20)
        lp.font.bold = True
        lp.font.color.rgb = ACCENT
        lp.font.name = "맑은 고딕"

    lb = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4.2), Inches(4.8))
    ltf = lb.text_frame
    ltf.word_wrap = True
    for i, item in enumerate(left_items):
        p = ltf.paragraphs[0] if i == 0 else ltf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_TEXT
        p.font.name = "맑은 고딕"
        p.space_after = Pt(8)

    # Right column
    if right_title:
        rt = slide.shapes.add_textbox(Inches(5.2), Inches(1.3), Inches(4.2), Inches(0.5))
        rp = rt.text_frame.paragraphs[0]
        rp.text = right_title
        rp.font.size = Pt(20)
        rp.font.bold = True
        rp.font.color.rgb = ACCENT
        rp.font.name = "맑은 고딕"

    rb = slide.shapes.add_textbox(Inches(5.2), Inches(1.8), Inches(4.2), Inches(4.8))
    rtf = rb.text_frame
    rtf.word_wrap = True
    for i, item in enumerate(right_items):
        p = rtf.paragraphs[0] if i == 0 else rtf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_TEXT
        p.font.name = "맑은 고딕"
        p.space_after = Pt(8)


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ── Slide 1: 표지 ──
    s1 = prs.slides.add_slide(blank)
    set_slide_bg(s1, PRIMARY)
    accent_bar = s1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(4.8), prs.slide_width, Inches(0.08)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = ACCENT
    accent_bar.line.fill.background()

    title_box = s1.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(1.5))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = "부동산 웹 플랫폼"
    tp.font.size = Pt(44)
    tp.font.bold = True
    tp.font.color.rgb = WHITE
    tp.font.name = "맑은 고딕"
    tp.alignment = PP_ALIGN.CENTER

    sub_box = s1.shapes.add_textbox(Inches(0.8), Inches(3.2), Inches(8.4), Inches(0.8))
    sp = sub_box.text_frame.paragraphs[0]
    sp.text = "제작 결과 발표"
    sp.font.size = Pt(32)
    sp.font.color.rgb = ACCENT
    sp.font.name = "맑은 고딕"
    sp.alignment = PP_ALIGN.CENTER

    info_box = s1.shapes.add_textbox(Inches(0.8), Inches(5.2), Inches(8.4), Inches(1.5))
    info_tf = info_box.text_frame
    for i, line in enumerate([
        "발표 시간: 20분  |  슬라이드: 10장",
        "팀: 김길동 · 홍길동 · 박미금 · 김서연",
        "2026년 7월",
    ]):
        p = info_tf.paragraphs[0] if i == 0 else info_tf.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
        p.font.color.rgb = LIGHT
        p.font.name = "맑은 고딕"
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(6)

    # ── Slide 2: 목차 ──
    s2 = prs.slides.add_slide(blank)
    set_slide_bg(s2, LIGHT)
    add_header_bar(s2, "목차 (Agenda)", prs)
    add_bullets(s2, [
        "1. 프로젝트 개요",
        "2. 팀원 소개",
        "3. 프로젝트 목표 및 범위",
        "4. 시스템 아키텍처",
        "5. 프론트엔드 기술 스택",
        "6. 백엔드 기술 스택",
        "7. RPA + AI (Python) 기술",
        "8. 주요 기능 및 구현 결과",
        "9. 배포 및 운영 환경",
        "10. 결론 및 Q&A",
    ], font_size=20)

    # ── Slide 3: 프로젝트 개요 ──
    s3 = prs.slides.add_slide(blank)
    set_slide_bg(s3, LIGHT)
    add_header_bar(s3, "1. 프로젝트 개요", prs)
    add_bullets(s3, [
        "프로젝트명: 부동산 웹 플랫폼",
        "목적: 매물 검색·등록·관리를 통합 제공하는 웹 서비스 구축",
        "대상 사용자: 일반 사용자(매물 탐색), 중개사/관리자(매물 관리)",
        "개발 기간: 기획 → 설계 → 개발 → 테스트 → 배포",
        "핵심 가치:",
        "   • 직관적인 UI/UX로 매물 정보 접근성 향상",
        "   • Spring Security 기반 안전한 회원·권한 관리",
        "   • RPA + AI로 부동산 데이터 수집·분석 자동화",
        "   • 클라우드 배포를 통한 안정적 서비스 운영",
    ], font_size=17)

    # ── Slide 4: 팀원 소개 ──
    s4 = prs.slides.add_slide(blank)
    set_slide_bg(s4, LIGHT)
    add_header_bar(s4, "2. 팀원 소개", prs)

    members = [
        ("김길동", "팀장 / 백엔드", "Spring Boot API, DB 설계, 보안"),
        ("홍길동", "프론트엔드", "HTML/CSS/JS, Bootstrap UI, 반응형 레이아웃"),
        ("박미금", "풀스택 / DevOps", "Gradle 빌드, 배포(TiDB·Render/EC2), CI/CD"),
        ("김서연", "RPA + AI", "Python 데이터 수집·분석, 시각화, ML 모델"),
    ]
    card_w, card_h = Inches(2.1), Inches(4.5)
    start_x = Inches(0.45)
    gap = Inches(0.25)
    for idx, (name, role, detail) in enumerate(members):
        x = start_x + idx * (card_w + gap)
        card = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.5), card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = ACCENT
        card.line.width = Pt(1.5)

        avatar = s4.shapes.add_shape(
            MSO_SHAPE.OVAL, x + Inches(0.65), Inches(1.8), Inches(0.8), Inches(0.8)
        )
        avatar.fill.solid()
        avatar.fill.fore_color.rgb = PRIMARY
        avatar.line.fill.background()
        ini = s4.shapes.add_textbox(x + Inches(0.65), Inches(1.95), Inches(0.8), Inches(0.6))
        ip = ini.text_frame.paragraphs[0]
        ip.text = name[0]
        ip.font.size = Pt(24)
        ip.font.bold = True
        ip.font.color.rgb = WHITE
        ip.font.name = "맑은 고딕"
        ip.alignment = PP_ALIGN.CENTER

        nb = s4.shapes.add_textbox(x + Inches(0.15), Inches(2.8), Inches(1.8), Inches(2.5))
        ntf = nb.text_frame
        ntf.word_wrap = True
        for j, (txt, bold, size, color) in enumerate([
            (name, True, 18, PRIMARY),
            (role, True, 14, ACCENT),
            (detail, False, 12, SUBTEXT),
        ]):
            p = ntf.paragraphs[0] if j == 0 else ntf.add_paragraph()
            p.text = txt
            p.font.bold = bold
            p.font.size = Pt(size)
            p.font.color.rgb = color
            p.font.name = "맑은 고딕"
            p.alignment = PP_ALIGN.CENTER
            p.space_after = Pt(8)

    # ── Slide 5: 목표 및 범위 ──
    s5 = prs.slides.add_slide(blank)
    set_slide_bg(s5, LIGHT)
    add_header_bar(s5, "3. 프로젝트 목표 및 범위", prs)
    add_two_column(
        s5,
        left_title="프로젝트 목표",
        left_items=[
            "매물 통합 검색·필터·상세 조회",
            "회원가입/로그인 및 권한별 접근 제어",
            "매물 등록·수정·삭제(CRUD) 관리",
            "관리자 대시보드 및 통계 제공",
            "부동산 시세·지역 데이터 AI 분석",
            "RPA 기반 외부 매물 데이터 자동 수집",
        ],
        right_title="개발 범위",
        right_items=[
            "웹 프론트: 반응형 UI (Bootstrap 5 + Flexbox)",
            "백엔드 API: Spring Boot REST + Thymeleaf",
            "DB: MySQL 8 / TiDB 연동",
            "인증: Spring Security + BCrypt",
            "빌드: Gradle 멀티 모듈 구성",
            "배포: Render 또는 AWS EC2",
        ],
    )

    # ── Slide 6: 시스템 아키텍처 ──
    s6 = prs.slides.add_slide(blank)
    set_slide_bg(s6, LIGHT)
    add_header_bar(s6, "4. 시스템 아키텍처", prs)

    layers = [
        ("Client Layer", "HTML · CSS · JS · Bootstrap 5 · Flexbox", PRIMARY),
        ("Application Layer", "Spring Boot 3.5 · Thymeleaf · Tomcat 10.1", ACCENT),
        ("Security Layer", "Spring Security · BCrypt · Session/JWT", RGBColor(0x8E, 0x44, 0xAD)),
        ("Data Layer", "MySQL 8 · JDBC Template · TiDB", RGBColor(0xE6, 0x7E, 0x22)),
        ("AI/RPA Layer", "Python · Requests · BeautifulSoup · scikit-learn", RGBColor(0x27, 0xAE, 0x60)),
        ("Deploy Layer", "Gradle Build · Render / AWS EC2", RGBColor(0xC0, 0x39, 0x2B)),
    ]
    for i, (layer, tech, color) in enumerate(layers):
        y = Inches(1.35 + i * 0.85)
        rect = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), y, Inches(7.6), Inches(0.65))
        rect.fill.solid()
        rect.fill.fore_color.rgb = color
        rect.line.fill.background()
        tb = s6.shapes.add_textbox(Inches(1.4), y + Inches(0.08), Inches(7.2), Inches(0.5))
        p = tb.text_frame.paragraphs[0]
        p.text = f"{layer}  →  {tech}"
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = "맑은 고딕"

    arrow_note = s6.shapes.add_textbox(Inches(1.2), Inches(6.5), Inches(7.6), Inches(0.5))
    ap = arrow_note.text_frame.paragraphs[0]
    ap.text = "▲ 3-Tier Architecture + RPA/AI Pipeline"
    ap.font.size = Pt(14)
    ap.font.color.rgb = SUBTEXT
    ap.font.name = "맑은 고딕"
    ap.alignment = PP_ALIGN.CENTER

    # ── Slide 7: 프론트엔드 ──
    s7 = prs.slides.add_slide(blank)
    set_slide_bg(s7, LIGHT)
    add_header_bar(s7, "5. 프론트엔드 기술 스택", prs)
    add_two_column(
        s7,
        left_title="핵심 기술",
        left_items=[
            "HTML5 — 시맨틱 마크업 구조",
            "CSS3 — 스타일링 및 애니메이션",
            "JavaScript (ES6+) — 동적 UI/UX",
            "Bootstrap 5 — 반응형 컴포넌트",
            "Flexbox Layout — 유연한 레이아웃",
        ],
        right_title="구현 포인트",
        right_items=[
            "모바일·태블릿·데스크톱 반응형 대응",
            "매물 카드 그리드 / 리스트 뷰 전환",
            "검색 필터 UI (지역·가격·면적)",
            "Thymeleaf 템플릿과 서버 렌더링 연동",
            "Bootstrap Navbar·Modal·Form 활용",
        ],
    )

    # ── Slide 8: 백엔드 ──
    s8 = prs.slides.add_slide(blank)
    set_slide_bg(s8, LIGHT)
    add_header_bar(s8, "6. 백엔드 기술 스택 & 의존성", prs)
    add_two_column(
        s8,
        left_title="런타임 & 프레임워크",
        left_items=[
            "JDK 21 (LTS)",
            "Spring Boot 3.5.x",
            "Thymeleaf (서버 사이드 템플릿)",
            "Tomcat 10.1 (내장 WAS)",
            "MySQL 8 (RDBMS)",
            "Gradle (빌드 도구)",
        ],
        right_title="Spring 의존성 (DI)",
        right_items=[
            "spring-boot-starter-web",
            "spring-boot-starter-security (BCrypt)",
            "spring-boot-starter-thymeleaf",
            "mysql-connector-j",
            "lombok",
            "spring-boot-starter-jdbc (JdbcTemplate)",
        ],
    )

    # ── Slide 9: RPA + AI ──
    s9 = prs.slides.add_slide(blank)
    set_slide_bg(s9, LIGHT)
    add_header_bar(s9, "7. RPA + AI (Python) 기술 스택", prs)
    add_two_column(
        s9,
        left_title="Python 패키지",
        left_items=[
            "NumPy — 수치 연산",
            "Pandas — 데이터 전처리·분석",
            "Matplotlib — 기본 시각화",
            "Seaborn — 통계 그래프",
            "Plotly — 인터랙티브 차트",
            "Requests — HTTP API/웹 요청",
            "BeautifulSoup — HTML 크롤링",
            "scikit-learn — ML 예측 모델",
        ],
        right_title="활용 시나리오",
        right_items=[
            "외부 부동산 사이트 매물 자동 수집 (RPA)",
            "지역별 시세·거래량 데이터 분석",
            "매물 가격 예측 모델 (회귀/분류)",
            "관리자 대시보드 시각화 리포트",
            "크롤링 데이터 → MySQL/TiDB 적재",
            "정기 배치 스케줄링 연동",
        ],
    )

    # ── Slide 10: 기능 + 배포 + 결론 ──
    s10 = prs.slides.add_slide(blank)
    set_slide_bg(s10, LIGHT)
    add_header_bar(s10, "8~10. 구현 결과 · 배포 · 결론", prs)

    sections = [
        ("주요 구현 기능", [
            "매물 검색·상세·등록·수정·삭제",
            "회원가입/로그인 (BCrypt 암호화)",
            "관리자 권한별 페이지 접근 제어",
            "RPA 데이터 수집 + AI 시세 분석",
        ]),
        ("배포 환경", [
            "DB: TiDB (MySQL 호환 분산 DB)",
            "서버: Render 또는 AWS EC2",
            "빌드: Gradle → JAR/WAR 배포",
            "HTTPS + 환경변수 기반 설정",
        ]),
        ("성과 & 향후", [
            "Full-Stack 부동산 플랫폼 완성",
            "팀 4인 역할 분담 협업 경험",
            "향후: 실시간 알림, 지도 API, 모바일 앱",
            "Q & A — 감사합니다!",
        ]),
    ]
    for i, (sec_title, items) in enumerate(sections):
        x = Inches(0.4 + i * 3.15)
        st = s10.shapes.add_textbox(x, Inches(1.3), Inches(2.9), Inches(0.45))
        sp2 = st.text_frame.paragraphs[0]
        sp2.text = sec_title
        sp2.font.size = Pt(17)
        sp2.font.bold = True
        sp2.font.color.rgb = ACCENT
        sp2.font.name = "맑은 고딕"

        sb = s10.shapes.add_textbox(x, Inches(1.8), Inches(2.9), Inches(4.8))
        stf = sb.text_frame
        stf.word_wrap = True
        for j, item in enumerate(items):
            p = stf.paragraphs[0] if j == 0 else stf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(13)
            p.font.color.rgb = DARK_TEXT
            p.font.name = "맑은 고딕"
            p.space_after = Pt(6)

    prs.save(OUTPUT)
    print(f"Created: {OUTPUT}")


if __name__ == "__main__":
    create_presentation()
